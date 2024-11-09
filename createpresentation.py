from pptx import Presentation

# Create a new PowerPoint presentation
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]

title.text = "Biliv-Augmenta Employee Attrition Analysis"
subtitle.text = "Insights and Recommendations\nPresented by: Frederick Kaakoye\nDate: 06/11/2024"

# Slide 2: Introduction
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title_2 = slide_2.shapes.title
content_2 = slide_2.placeholders[1]

title_2.text = "Introduction"
content_2.text = (
    "Objective: Understand and reduce employee attrition at Biliv-Augmenta.\n"
    "Scope: Analysis of demographics, performance, and attrition trends using Power BI."
)

# Slide 3: Overview of Attrition Metrics
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title_3 = slide_3.shapes.title
content_3 = slide_3.placeholders[1]

title_3.text = "Overview of Attrition Metrics"
content_3.text = (
    "Key Metrics:\n"
    "- Total Employees\n"
    "- Active vs. Inactive Employees\n"
    "- Attrition Rate\n"
    "Visual: KPI Cards displaying each metric."
)

# Slide 4: Attrition Trends Over Time
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title_4 = slide_4.shapes.title
content_4 = slide_4.placeholders[1]

title_4.text = "Attrition Trends Over Time"
content_4.text = (
    "Insight: Annual attrition rates and identified spikes.\n"
    "Visual: Line Chart showing Attrition Rate by Year."
)

# Slide 5: Demographics Analysis
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title_5 = slide_5.shapes.title
content_5 = slide_5.placeholders[1]

title_5.text = "Demographics Analysis"
content_5.text = (
    "Key Findings:\n"
    "- Higher attrition in younger employees (<2 years tenure).\n"
    "- Gender disparities in attrition rates.\n"
    "- Marital status correlation with retention.\n"
    "Visuals: Stacked Column and Donut Charts."
)

# Slide 6: Performance Tracker Insights
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title_6 = slide_6.shapes.title
content_6 = slide_6.placeholders[1]

title_6.text = "Performance Tracker Insights"
content_6.text = (
    "Key Findings:\n"
    "- Regular review cycles correlate with lower attrition.\n"
    "- Low self and manager ratings linked to higher turnover.\n"
    "- Career development opportunities matter.\n"
    "Visuals: Line Graph, Bar Chart."
)

# Slide 7: Attrition Drivers
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title_7 = slide_7.shapes.title
content_7 = slide_7.placeholders[1]

title_7.text = "Attrition Drivers"
content_7.text = (
    "Primary Drivers:\n"
    "- Limited career advancement.\n"
    "- Workload stress and overtime.\n"
    "- Compensation dissatisfaction.\n"
    "Visual: Pie Chart or Bar Graph highlighting reasons for attrition."
)

# Slide 8: Departmental and Role-Based Trends
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title_8 = slide_8.shapes.title
content_8 = slide_8.placeholders[1]

title_8.text = "Departmental and Role-Based Trends"
content_8.text = (
    "Key Findings:\n"
    "- High attrition in specific departments (e.g., roles requiring extensive travel).\n"
    "- Correlation between tenure in role and turnover.\n"
    "Visuals: Stacked Column Chart, Heat Map."
)

# Slide 9: Compensation and Benefits Analysis
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
title_9 = slide_9.shapes.title
content_9 = slide_9.placeholders[1]

title_9.text = "Compensation and Benefits Analysis"
content_9.text = (
    "Key Findings:\n"
    "- Higher average salaries linked to lower attrition.\n"
    "- Impact of salary hikes and stock options on retention.\n"
    "Visual: Scatter Plot showing Monthly Income vs. Attrition Rate."
)

# Slide 10: Geographic and Commute Factors
slide_10 = prs.slides.add_slide(prs.slide_layouts[1])
title_10 = slide_10.shapes.title
content_10 = slide_10.placeholders[1]

title_10.text = "Geographic and Commute Factors"
content_10.text = (
    "Key Findings:\n"
    "- Longer commute distances associated with higher attrition.\n"
    "- Regional trends influencing retention.\n"
    "Visual: Map Visualization of Attrition Rate by State/Region."
)

# Slide 11: Work-Life Balance and Overtime Impact
slide_11 = prs.slides.add_slide(prs.slide_layouts[1])
title_11 = slide_11.shapes.title
content_11 = slide_11.placeholders[1]

title_11.text = "Work-Life Balance and Overtime Impact"
content_11.text = (
    "Key Findings:\n"
    "- High overtime correlates with increased turnover.\n"
    "- Work-life balance satisfaction influences retention.\n"
    "Visuals: Bar Chart showing Overtime vs. Attrition Rate."
)

# Slide 12: Recommendations
slide_12 = prs.slides.add_slide(prs.slide_layouts[1])
title_12 = slide_12.shapes.title
content_12 = slide_12.placeholders[1]

title_12.text = "Recommendations"
content_12.text = (
    "Strategies to Reduce Attrition:\n"
    "- Enhance career development programs.\n"
    "- Improve compensation and benefits packages.\n"
    "- Implement flexible work arrangements.\n"
    "- Strengthen managerial support and regular feedback mechanisms."
)

# Slide 13: Action Plan
slide_13 = prs.slides.add_slide(prs.slide_layouts[1])
title_13 = slide_13.shapes.title
content_13 = slide_13.placeholders[1]

title_13.text = "Action Plan"
content_13.text = (
    "Short-Term Actions:\n"
    "- Schedule regular performance reviews.\n"
    "- Conduct salary benchmarking.\n\n"
    "Long-Term Actions:\n"
    "- Develop clear career progression paths.\n"
    "- Introduce wellness and work-life balance initiatives."
)

# Slide 14: Conclusion
slide_14 = prs.slides.add_slide(prs.slide_layouts[1])
title_14 = slide_14.shapes.title
content_14 = slide_14.placeholders[1]

title_14.text = "Conclusion"
content_14.text = (
    "Summary of Insights:\n"
    "- Recap key findings and their implications.\n\n"
    "Next Steps:\n"
    "- Outline immediate actions and future analysis.\n"
)

# Slide 15: Thank You
slide_15 = prs.slides.add_slide(prs.slide_layouts[1])
title_15 = slide_15.shapes.title
content_15 = slide_15.placeholders[1]

title_15.text = "Thank You"
content_15.text = "Thank you for your attention!\n\nContact Information: Frederick Kaakoye | Email: frederickkaakoye@gmail.com | LinkedIn: https://www.linkedin.com/in/frederick-kaakoye-98ab16243/"

# Save the presentation
pptx_file_path = "Biliv-Augmenta_Employee_Attrition_Analysis_Report.pptx"
prs.save(pptx_file_path)

print(f"Presentation saved as {pptx_file_path}")